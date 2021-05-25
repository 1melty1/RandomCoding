import pptx
from pptx import presentation


print("plz save the ppt before use")

ppt_path = input("Enter the path of the ppt(Drag and drop): ")
txt_path = input("Enter the path of the lyrics file(Drag and drop): ")
song_num = int(input("How many songs in total"))
lower = int(input("Starting page (Inclusive): "))
upper = int(input("Ending page (Inclusive): "))
song_list = []

def cut_N_replace(lower, upper, ppt_path, txt_path, mode):
    f = open(txt_path, "r")
    ppt = presentation(ppt_path)
    current_sentence = ""

    for i in range(lower, upper):
        slide = get(i)
        song_name = ""
        for j in range(lower, upper+1):
            to_be_added = f.readln()
            if current_sentence[1] == "#":
                song_name = to_be_added
            else:
                current_sentence = current_sentence + '\n' + to_be_added

            if current_sentence[-1] == "#": #if it's the end of sentence                
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text == "1":
                        shape.text = current_sentence
                    if shape.has_text_frame and shape.text == "2":
                        shape.text = song_name        

                current_sentence == "" #reset the current sentence

    f.close()

for i in range(song_num):
    song_list[i] = input("Song name of song ", i)
    Mode = input("Select mode \n(1.Auto(currently not supported) \n2.Manual) \nEnter mode: ")
    cut_N_replace(lower, upper, ppt_path, txt_path, mode)

while confirm != "y":
    confirm = input("Enter Y to confirm... \n")
